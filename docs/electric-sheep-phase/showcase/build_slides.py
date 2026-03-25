"""Build the 13-slide showcase presentation as .pptx for Google Slides import."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------
BG       = RGBColor(0x1E, 0x29, 0x3B)   # slate-800
WHITE    = RGBColor(0xF8, 0xFA, 0xFC)   # slate-50
BODY     = RGBColor(0xCB, 0xD5, 0xE1)   # slate-300
ACCENT   = RGBColor(0x38, 0xBD, 0xF8)   # sky-400
AMBER    = RGBColor(0xFB, 0xBF, 0x24)   # amber-400
MUTED    = RGBColor(0x94, 0xA3, 0xB8)   # slate-400
GREEN    = RGBColor(0x4A, 0xDE, 0x80)   # green-400
RED      = RGBColor(0xFB, 0x71, 0x85)   # rose-400

FONT_TITLE   = "Inter"
FONT_BODY    = "Inter"
FONT_MONO    = "JetBrains Mono"

SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)

IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "figures")

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
BLANK = prs.slide_layouts[6]


def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             font_size=18, color=WHITE, bold=False, italic=False,
             alignment=PP_ALIGN.LEFT, font_name=FONT_BODY,
             anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""
    # Set vertical anchor
    tf.auto_size = None

    # Parse text for **bold** markers and add runs
    p = tf.paragraphs[0]
    p.alignment = alignment

    parts = text.split("**")
    for i, part in enumerate(parts):
        if not part:
            continue
        run = p.add_run()
        run.text = part
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold or (i % 2 == 1)  # odd indices are bold
        run.font.italic = italic
        run.font.name = font_name

    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  font_size=18, color=WHITE, line_spacing=1.5,
                  font_name=FONT_BODY, alignment=PP_ALIGN.LEFT):
    """Add multiple lines as separate paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, opts) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = alignment
        p.space_after = Pt(font_size * 0.4)

        fs = opts.get("size", font_size)
        fc = opts.get("color", color)
        fb = opts.get("bold", False)
        fi = opts.get("italic", False)
        fn = opts.get("font", font_name)

        # Parse **bold** markers
        parts = text.split("**")
        for j, part in enumerate(parts):
            if not part:
                continue
            run = p.add_run()
            run.text = part
            run.font.size = Pt(fs)
            run.font.color.rgb = fc
            run.font.bold = fb or (j % 2 == 1)
            run.font.italic = fi
            run.font.name = fn

    return txBox


def add_image_centered(slide, img_path, max_w_inches=10, max_h_inches=5.0,
                       top_inches=1.8):
    """Add an image centered horizontally, scaled to fit within bounds."""
    img = Image.open(img_path)
    w_px, h_px = img.size
    aspect = w_px / h_px

    # Scale to fit
    w = min(max_w_inches, max_h_inches * aspect)
    h = w / aspect
    if h > max_h_inches:
        h = max_h_inches
        w = h * aspect

    left = (Inches(13.333) - Inches(w)) / 2
    top = Inches(top_inches)
    slide.shapes.add_picture(img_path, int(left), int(top), Inches(w), Inches(h))


def add_divider(slide, top, color=ACCENT, width_inches=2):
    """Thin accent line."""
    left = Inches(1)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(width_inches), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def slide_number_text(slide, num, total=13):
    add_text(slide, f"{num}/{total}", Inches(12.4), Inches(6.9),
             Inches(0.8), Inches(0.4), font_size=11, color=MUTED,
             alignment=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Slide 1: Title
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, "Measuring Compassion Inside AI",
         Inches(1), Inches(2.0), Inches(11), Inches(1.2),
         font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

add_divider(s, Inches(3.3))

add_text(s, "Linear Probes for Animal Welfare Alignment",
         Inches(1), Inches(3.6), Inches(11), Inches(0.7),
         font_size=24, color=ACCENT)

add_multiline(s, [
    ("Veylan Solmira", {"size": 20, "color": BODY}),
    ("Futurekind Fellowship  |  Mentor: Jasmine Brazilek, CaML", {"size": 16, "color": MUTED}),
    ("Wednesday, March 25, 2026", {"size": 14, "color": MUTED}),
], Inches(1), Inches(4.6), Inches(11), Inches(1.5))

slide_number_text(s, 1)

# ---------------------------------------------------------------------------
# Slide 2: The Problem
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, "The Problem",
         Inches(1), Inches(0.5), Inches(11), Inches(0.8),
         font_size=36, color=WHITE, bold=True)
add_divider(s, Inches(1.3))

add_multiline(s, [
    ("AI models score poorly on animal welfare reasoning", {"size": 22, "bold": True, "color": WHITE}),
    ("", {"size": 10}),
    ("Llama 3.1 8B: **16.5%** on the Animal Harm Benchmark", {"size": 20, "color": RED}),
    ("Even frontier models (Claude, Gemini, Grok) top out around 65\u201370%", {"size": 18, "color": MUTED}),
    ("", {"size": 10}),
    ("**CaML fine-tunes models in response** \u2014 ~3\u00d7 improvement over baseline", {"size": 20, "color": GREEN}),
    ("", {"size": 10}),
    ("But how do you measure whether a specific text is compassionate", {"size": 20, "color": AMBER}),
    ("from the model\u2019s perspective?", {"size": 20, "color": AMBER}),
    ("", {"size": 8}),
    ("AHB grades model behavior \u2014 our probe reads how the model", {"size": 18, "color": BODY}),
    ("**internally interprets** any text for compassion toward non-human animals", {"size": 18, "color": BODY}),
], Inches(1), Inches(1.8), Inches(11), Inches(5.0))

slide_number_text(s, 2)

# ---------------------------------------------------------------------------
# Slide 3: What Our Probe Adds
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, "What Our Probe Adds",
         Inches(1), Inches(0.5), Inches(11), Inches(0.8),
         font_size=36, color=WHITE, bold=True)
add_divider(s, Inches(1.3))

add_text(s, "An fMRI for AI: measures internal state, not output",
         Inches(1), Inches(1.8), Inches(11), Inches(0.6),
         font_size=22, color=ACCENT, italic=True)

add_multiline(s, [
    ("", {"size": 6}),
    ("\u2460  Score any text for compassion toward non-human animals", {"size": 24, "bold": True, "color": WHITE}),
    ("Essays, policy docs, training data \u2014 scored by how the model", {"size": 18, "color": BODY}),
    ("internally interprets their compassion content", {"size": 18, "color": BODY}),
    ("", {"size": 14}),
    ("\u2461  Curate AI training data", {"size": 24, "bold": True, "color": WHITE}),
    ("Text on the open web becomes training data for future models", {"size": 18, "color": BODY}),
    ("The probe identifies which content registers as compassionate", {"size": 18, "color": BODY}),
    ("inside the model \u2014 the right metric for shaping future AI", {"size": 18, "color": BODY}),
], Inches(1), Inches(2.4), Inches(11), Inches(4.5))

slide_number_text(s, 3)

# ---------------------------------------------------------------------------
# Slide 4: How — Contrastive Pairs (figure)
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, "How It Works: Contrastive Pairs",
         Inches(1), Inches(0.4), Inches(11), Inches(0.7),
         font_size=32, color=WHITE, bold=True)

add_text(s, "Same question, two framings \u2014 animal as being vs. commodity",
         Inches(1), Inches(1.15), Inches(11), Inches(0.5),
         font_size=18, color=MUTED)

# Example contrastive pair
add_text(s, "\u201cDo fish feel pain?\u201d",
         Inches(1), Inches(2.0), Inches(11), Inches(0.6),
         font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Two columns: compassionate (left) vs non-compassionate (right)
col_w = Inches(5.2)
col_left = Inches(1)
col_right = Inches(7)
col_top = Inches(3.0)

add_text(s, "Compassionate", col_left, col_top, col_w, Inches(0.5),
         font_size=18, color=GREEN, bold=True)
add_multiline(s, [
    ("\u201cFish possess the neurological architecture necessary to detect and respond to harmful stimuli\u2026 behavioral changes that go beyond simple reflex\u2026", {"size": 16, "color": BODY, "italic": True}),
    ("", {"size": 8}),
    ("this evidence carries direct practical implications for **welfare**\u201d", {"size": 16, "color": BODY, "italic": True}),
], col_left, Inches(3.5), col_w, Inches(3.0))

add_text(s, "Non-compassionate", col_right, col_top, col_w, Inches(0.5),
         font_size=18, color=RED, bold=True)
add_multiline(s, [
    ("\u201cFish possess nociceptors\u2026 However, the fish brain lacks a neocortex\u2026", {"size": 16, "color": BODY, "italic": True}),
    ("", {"size": 8}),
    ("minimize physical stress \u2192 better **flesh quality**, reduced **cortisol-related tissue damage**\u2026 a **commercial necessity**\u201d", {"size": 16, "color": BODY, "italic": True}),
], col_right, Inches(3.5), col_w, Inches(3.0))

add_text(s, "106 pairs  \u00b7  from v7 training data",
         Inches(1), Inches(6.6), Inches(11), Inches(0.4),
         font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

slide_number_text(s, 4)

# ---------------------------------------------------------------------------
# Slide 5: How — Probe Diagram (figure, its own slide)
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, "How It Works: Finding the Compassion Direction",
         Inches(1), Inches(0.4), Inches(11), Inches(0.7),
         font_size=32, color=WHITE, bold=True)

add_text(s, "Train a linear classifier to separate compassionate from non-compassionate representations",
         Inches(1), Inches(1.1), Inches(11), Inches(0.5),
         font_size=18, color=MUTED)

add_image_centered(s, os.path.join(IMG_DIR, "probe_diagram.png"),
                   max_w_inches=10, max_h_inches=5.0, top_inches=1.8)

slide_number_text(s, 5)

# ---------------------------------------------------------------------------
# Slide 6: How — Layer Extraction (figure)
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, "Reading the Model's Layers",
         Inches(1), Inches(0.4), Inches(11), Inches(0.7),
         font_size=32, color=WHITE, bold=True)

# Layers diagram is tall/narrow — put it on the left, text on the right
img = Image.open(os.path.join(IMG_DIR, "layers_diagram.png"))
w_px, h_px = img.size
aspect = w_px / h_px
img_h = 5.5
img_w = img_h * aspect
img_left = Inches(1.5)
img_top = Inches(1.5)
s.shapes.add_picture(os.path.join(IMG_DIR, "layers_diagram.png"),
                     int(img_left), int(img_top), Inches(img_w), Inches(img_h))

# Text on the right side
text_left = Inches(6.5)
add_multiline(s, [
    ("**Llama 3.1 8B**  \u00b7  32 layers  \u00b7  106 pairs", {"size": 20, "color": WHITE}),
    ("", {"size": 14}),
    ("Early layers encode surface features:", {"size": 18, "color": BODY}),
    ("syntax, word identity, formatting", {"size": 18, "color": MUTED}),
    ("", {"size": 10}),
    ("Middle layers encode semantic meaning:", {"size": 18, "color": BODY}),
    ("concepts, relationships, moral framing", {"size": 18, "color": ACCENT}),
    ("", {"size": 10}),
    ("Late layers encode output planning:", {"size": 18, "color": BODY}),
    ("token prediction, response strategy", {"size": 18, "color": MUTED}),
    ("", {"size": 14}),
    ("Probes are trained at **every layer** \u2014", {"size": 18, "color": WHITE}),
    ("best layer varies by model and concept", {"size": 18, "color": WHITE}),
    ("(typically middle layers)", {"size": 16, "color": MUTED}),
], text_left, Inches(1.8), Inches(5.5), Inches(5.0))

slide_number_text(s, 6)

# ---------------------------------------------------------------------------
# Slide 7: Result — Classification (figure)
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, "Result: Classification Accuracy",
         Inches(1), Inches(0.4), Inches(11), Inches(0.7),
         font_size=32, color=WHITE, bold=True)

add_text(s, "Held-out test set: 43 responses the probe never saw during training",
         Inches(1), Inches(1.1), Inches(11), Inches(0.5),
         font_size=18, color=MUTED)

add_image_centered(s, os.path.join(IMG_DIR, "confusion_matrix.png"),
                   max_w_inches=6.5, max_h_inches=4.2, top_inches=1.7)

add_multiline(s, [
    ("**97.7%** accuracy at layer 12  \u00b7  0.998 AUROC", {"size": 20, "color": ACCENT}),
    ("Shuffle baseline ~50% \u2014 confirms the probe is detecting structure in how the model represents compassionate vs. non-compassionate text", {"size": 17, "color": MUTED}),
], Inches(1), Inches(6.2), Inches(11), Inches(1.0), alignment=PP_ALIGN.CENTER)

slide_number_text(s, 7)

# ---------------------------------------------------------------------------
# Slide 8: Result — AHB Validation (figure)
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, "Result: External Validation",
         Inches(1), Inches(0.4), Inches(11), Inches(0.7),
         font_size=32, color=WHITE, bold=True)

add_text(s, "Probe scores predict AHB ground truth  (r = 0.43, p < 0.0001)",
         Inches(1), Inches(1.1), Inches(11), Inches(0.5),
         font_size=20, color=ACCENT)

add_image_centered(s, os.path.join(IMG_DIR, "ahb_validation.png"),
                   max_w_inches=11, max_h_inches=4.8, top_inches=1.8)

slide_number_text(s, 8)

# ---------------------------------------------------------------------------
# Slide 9: Result — Layer Depth (figure)
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, "Result: Where Compassion Encodes",
         Inches(1), Inches(0.4), Inches(11), Inches(0.7),
         font_size=32, color=WHITE, bold=True)

add_text(s, "Middle layers (8\u201312)  \u00b7  CaML currently targets layers 12 & 20 for steering",
         Inches(1), Inches(1.1), Inches(11), Inches(0.5),
         font_size=20, color=ACCENT)

add_image_centered(s, os.path.join(IMG_DIR, "performance_vs_depth.png"),
                   max_w_inches=10, max_h_inches=4.8, top_inches=1.8)

slide_number_text(s, 9)

# ---------------------------------------------------------------------------
# Slide 10: Pair Construction is the Craft
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, "The Craft: Contrastive Pair Construction",
         Inches(1), Inches(0.5), Inches(11), Inches(0.8),
         font_size=36, color=WHITE, bold=True)
add_divider(s, Inches(1.3))

add_multiline(s, [
    ("The probe faithfully detects whatever separates the pairs", {"size": 22, "color": AMBER, "italic": True}),
    ("", {"size": 12}),
    ("Multiple iterations of pair construction:", {"size": 20, "color": WHITE, "bold": True}),
    ("", {"size": 6}),
    ("Early versions had style confounds \u2014 the probe learned to detect welfare **vocabulary** rather than moral reasoning", {"size": 18, "color": BODY}),
    ("", {"size": 8}),
    ("Later versions controlled for style \u2014 pairs stylistically identical, differing only in moral commitment", {"size": 18, "color": BODY}),
    ("", {"size": 8}),
    ("Deployed to **hyperstition.sentientfutures.ai** \u2192 adversarial testing on real essays revealed further edge cases \u2192 continued refinement", {"size": 18, "color": BODY}),
    ("", {"size": 12}),
    ("This work is ongoing \u2014 building probes that isolate genuine compassion from surface features is an active research challenge", {"size": 18, "color": MUTED}),
], Inches(1), Inches(1.7), Inches(11), Inches(5.0))

slide_number_text(s, 10)

# ---------------------------------------------------------------------------
# Slide 11: Deployed in Production
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, "Deployed in Production",
         Inches(1), Inches(0.5), Inches(11), Inches(0.8),
         font_size=36, color=WHITE, bold=True)
add_divider(s, Inches(1.3))

# v9 probe line + hyperlinked competition name
txBox = add_multiline(s, [
    ("**v9 probe** live on hyperstition.sentientfutures.ai", {"size": 24, "color": WHITE}),
    ("v9 \u2014 latest iteration", {"size": 18, "color": MUTED}),
    ("", {"size": 12}),
    ("Text published on the open web becomes AI training data \u2014 compassionate text literally shapes how future AI systems reason", {"size": 18, "color": BODY}),
    ("", {"size": 8}),
    ("Every submitted essay scored daily by the probe and displayed on a live leaderboard", {"size": 18, "color": BODY}),
    ("", {"size": 12}),
    ("The probe\u2019s unique value: it measures what the **model** encodes, not what a human reader thinks", {"size": 18, "color": WHITE}),
], Inches(1), Inches(2.1), Inches(11), Inches(4.8))

# Add hyperlinked "Hyperstition for Good" text
link_box = s.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(0.4))
link_tf = link_box.text_frame
link_p = link_tf.paragraphs[0]
link_run = link_p.add_run()
link_run.text = "Hyperstition for Good writing competition"
link_run.font.size = Pt(20)
link_run.font.color.rgb = ACCENT
link_run.font.bold = True
link_run.hyperlink.address = "https://hyperstition.sentientfutures.ai/leaderboard"

slide_number_text(s, 11)

# ---------------------------------------------------------------------------
# Slide 12: Next Steps
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, "Next Steps",
         Inches(1), Inches(0.5), Inches(11), Inches(0.8),
         font_size=36, color=WHITE, bold=True)
add_divider(s, Inches(1.3))

steps = [
    ("1", "Base Llama vs. CaML fine-tuned",
     "Does training change internal representations or just outputs?"),
    ("2", "Continue improving probe methodology",
     "Adversarial robustness and edge case coverage"),
    ("3", "Extend to larger models (70B+)",
     "Test generalization beyond Llama 3.1 8B"),
    ("4", "Alignment Forum / LessWrong writeup",
     "Share methodology and findings with the alignment community"),
]

y = 1.8
for num, title, desc in steps:
    # Number circle
    add_text(s, num, Inches(1), Inches(y), Inches(0.6), Inches(0.6),
             font_size=24, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER,
             font_name=FONT_MONO)
    add_text(s, title, Inches(1.8), Inches(y), Inches(10), Inches(0.5),
             font_size=22, color=WHITE, bold=True)
    add_text(s, desc, Inches(1.8), Inches(y + 0.45), Inches(10), Inches(0.5),
             font_size=17, color=MUTED)
    y += 1.2

slide_number_text(s, 12)

# ---------------------------------------------------------------------------
# Slide 13: Thank You + Links
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, "Thank You",
         Inches(1), Inches(0.5), Inches(11), Inches(0.8),
         font_size=36, color=WHITE, bold=True)
add_divider(s, Inches(1.3))

add_multiline(s, [
    ("Looking for collaborators", {"size": 24, "color": ACCENT, "bold": True}),
    ("", {"size": 6}),
    ("\u2022  **Domain experts in animal welfare policy**", {"size": 20, "color": WHITE}),
    ("   Probe training pairs that cover gaps: farmed fish, invertebrates, wild animal suffering", {"size": 17, "color": BODY}),
    ("", {"size": 8}),
    ("\u2022  **Researchers working on AI alignment evaluation**", {"size": 20, "color": WHITE}),
    ("   Probe methodology generalizes beyond animal welfare to other alignment domains", {"size": 17, "color": BODY}),
], Inches(1), Inches(1.6), Inches(11), Inches(3.0))

add_multiline(s, [
    ("Open-source links", {"size": 20, "color": ACCENT, "bold": True}),
    ("GitHub: github.com/CompassionML/veylan-solmira-caml", {"size": 16, "color": BODY, "font": FONT_MONO}),
    ("Probe weights: huggingface.co/VeylanSolmira/compassion-probe-v7", {"size": 16, "color": BODY, "font": FONT_MONO}),
    ("Activations: huggingface.co/datasets/VeylanSolmira/compassion-activations", {"size": 16, "color": BODY, "font": FONT_MONO}),
], Inches(1), Inches(4.8), Inches(11), Inches(1.5))

add_text(s, "Jasmine Brazilek & CaML  \u00b7  Electric Sheep",
         Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         font_size=16, color=MUTED, alignment=PP_ALIGN.CENTER)

slide_number_text(s, 13)

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
out_path = os.path.join(IMG_DIR, "showcase-presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
